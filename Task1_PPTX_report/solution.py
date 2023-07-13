# Szekeres Ádám solution

# imports
import json, csv, time, collections, collections.abc, logging
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

# getting .json file name, if not found, ask again
while (1):
    json_filename=input("Please type the .json filename! (eg. sample) ")
    try:
        json_file=open(json_filename+".json",'r')
        break
    except FileNotFoundError:
        print("File not found!")

json_content=json.loads(json_file.read()) # load file content into json object
json_file.close()

no_error=True
prs = Presentation()
FORMAT = '%(asctime)s %(levelname)s %(message)s'
logging.basicConfig(format=FORMAT, filename="Logfile.log", level=logging.INFO)

for current_slide in json_content["presentation"]: #looping through slides in json one by one, the "type" determines which slide is added

    if current_slide["type"]=="title":
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = current_slide["title"]
        slide.placeholders[1].text = current_slide["content"]
        
    elif current_slide["type"]=="text":
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = current_slide["title"]
        txBox = slide.shapes.add_textbox(Cm(3.5), Cm(3), Cm(14.5), Cm(30))
        txBox.text_frame.text=current_slide["content"]
        
    elif current_slide["type"]=="list":
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = current_slide["title"]
        for line in current_slide["content"]:
            p=slide.shapes.placeholders[1].text_frame.add_paragraph()
            p.text=line["text"]
            p.level=line["level"]
            
    elif current_slide["type"]=="picture":
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = current_slide["title"]
        try: # error handling if no picture is found with the given name
            slide.shapes.add_picture(current_slide["content"],Cm(3.5), Cm(3))
        except FileNotFoundError:
            print(current_slide["content"]," not found! Exiting without creating pptx.")
            logging.error(current_slide["content"]+" not found! pptx not created.")
            no_error=False
            break
        
    elif current_slide["type"]=="plot":
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = current_slide["title"]

        chart_data = ChartData()
        series_1 = chart_data.add_series("")
        category_1 = chart_data.add_category("")
        try: # error handling if no csv (dat) file is found with the given name
            csv_file = open(current_slide["content"],"r")
        except FileNotFoundError:
            print(current_slide["content"]," not found! Exiting without creating pptx.")
            logging.error(current_slide["content"]+" not found! pptx not created.")
            no_error=False
            break
        csv_reader=csv.reader(csv_file, delimiter=';')
        for data_points in csv_reader:
            series_1.add_data_point(data_points[0],data_points[1])
        csv_file.close()
        chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE, Cm(3.5), Cm(3), Inches(6.4), Inches(4.8), chart_data).chart
        chart.value_axis.axis_title.text_frame.text = current_slide["configuration"]["y-label"]
        chart.category_axis.axis_title.text_frame.text = current_slide["configuration"]["x-label"]
        chart.has_legend=False

    logging.info(current_slide["type"]+" slide created")

if no_error: # if there were no errors, save the file and exit
    pptx_filename=input("Type a filename for the pptx (eg. output): ")
    prs.save(pptx_filename+".pptx")
    print("Output saved in "+pptx_filename+".pptx")
    logging.info("Output saved in "+pptx_filename+".pptx")
    print("Exiting...")
time.sleep(2)
